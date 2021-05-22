# Made by CrumblyLiquid in 2021

import os
from pathlib import Path
from pptx import Presentation
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt

# Go through every .pptx file in that directory
for filepath in os.listdir(Path(__file__).parent.absolute()):
    if filepath.endswith('.pptx'):
        d = Document()
        title_separator = "<title-separator>"
        p = Presentation(filepath)
        text = []
        previous = ""
        skips = ["Zdroje", "Zdroje:", "Obsah", "Obsah:", "Doplňovačka", "Doplňovačka:"]
        for slide in p.slides:
            title = None
            parts = []
            for shape in slide.shapes:
                # Only shapes with text
                if hasattr(shape, "text"):
                    # Check for title shape
                    if (shape == slide.shapes.title and shape.text == slide.shapes.title.text):
                        t = shape.text.replace("\x0b", "")
                        t = t.rstrip(" ").rstrip("\t")
                        if (t != ""):
                            title = t
                    # For every other shape
                    else:
                        for paragraph in shape.text_frame.paragraphs:
                            t = paragraph.text.replace("\x0b", "")
                            t = t.rstrip(" ").rstrip("\t")
                            # Only text with something in it
                            if (t != ""):
                                parts.append(t)
            # For first slide only
            if (p.slides.index(slide) == 0 and title is not None):
                dtitle = d.add_paragraph()
                dtitle.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.CENTER
                run = dtitle.add_run(title)
                font = run.font
                font.name = 'Calibri'
                font.bold = True
                font.size = Pt(20)
            # Skip slide
            elif (title in skips):
                continue
            elif (len(parts) > 0):
                if (title is not None):
                    # Add title if it's a different than the last one
                    if (title != previous):
                        text.append(title_separator+title)
                    # Remove \n if both slides have the same title
                    else:
                        # We remove the last item of the list which should be a ""
                        # thus removing the \n that would othervise be there
                        text.pop(-1)
                    previous = title
                # Add paragraphs
                text.extend(parts)
                # Add slide splitter
                text.append("") # Will turn into \n

        for t in text:
            if (t.startswith(title_separator)):
                paragraph = d.add_paragraph()
                run = paragraph.add_run(t.replace(title_separator, ""))
                font = run.font
                font.name = 'Calibri'
                font.size = Pt(16)
            elif (t == ""):
                d.add_paragraph(t)
            else:
                paragraph = d.add_paragraph(style='List Bullet')
                run = paragraph.add_run(t)
                font = run.font
                font.name = 'Calibri'
                font.size = Pt(12)

        # Save .docx file
        txt_name = filepath.split("/")[-1].replace(".pptx", "") + ".docx"
        txt_path = str(Path(__file__).parent.absolute() / txt_name)
        d.save(txt_path)